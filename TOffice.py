import os
import re
import zipfile
from distutils.dir_util import copy_tree
import shutil
import pandas as pd
import win32com.client

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document

from fpdf import FPDF
from PIL import Image

from PyPDF2 import PdfReader, PdfWriter


# 파일 내 이미지 한번에 저장하기 (결과 저장 폴더, 작업 대상 폴더) 
def D_Img(TPath, result_f):

    # 결과를 저장할 폴더와 임시 저장 폴더 이름을 변수로 지정
    temp_f = 'Temp'

    # 작업할 폴더 선택하고 파워포인트, 엑셀, 워드 파일별로 리스트 만들기
    pptfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]ppt', f)]
    xlfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]xls', f)]
    wordfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]doc', f)]

    # PPT 파일 반복문
    for pptfile in pptfiles:

        # 저장할 파일 경로 생성
        newpath = os.path.join(result_f, 'ppt', pptfile[:pptfile.find('.')]) 
        if not os.path.exists(newpath):
            os.makedirs(newpath)

        # 임시 폴더에 MS Office 파일 압축 해제
        with zipfile.ZipFile(os.path.join(TPath,pptfile), 'r') as zip_ref:
            zip_ref.extractall(temp_f)

        # 임시 폴더에서 이미지 파일을 저장할 파일 경로에 복사
        copy_tree(os.path.join(temp_f,'ppt','media'), newpath)

        # 임시 폴더 삭제
        shutil.rmtree(temp_f)

    # Excel 파일 반복문
    for xlfile in xlfiles:

        # 저장할 파일 경로 생성
        newpath = os.path.join(result_f, 'xl', xlfile[:xlfile.find('.')]) 
        if not os.path.exists(newpath):
            os.makedirs(newpath)

        # 임시 폴더에 MS Office 파일 압축 해제
        with zipfile.ZipFile(os.path.join(TPath,xlfile), 'r') as zip_ref:
            zip_ref.extractall(temp_f)

        # 임시 폴더에서 이미지 파일을 저장할 파일 경로에 복사
        copy_tree(os.path.join(temp_f,'xl','media'), newpath)

        # 임시 폴더 삭제
        shutil.rmtree(temp_f)

    # Word 파일 반복문
    for wordfile in wordfiles:

        # 저장할 파일 경로 생성
        newpath = os.path.join(result_f, 'word', wordfile[:wordfile.find('.')]) 
        if not os.path.exists(newpath):
            os.makedirs(newpath)

        # 임시 폴더에 MS Office 파일 압축 해제
        with zipfile.ZipFile(os.path.join(TPath,wordfile), 'r') as zip_ref:
            zip_ref.extractall(temp_f)

        # 임시 폴더에서 이미지 파일을 저장할 파일 경로에 복사
        copy_tree(os.path.join(temp_f,'word','media'), newpath)

        # 임시 폴더 삭제
        shutil.rmtree(temp_f)

# 같은 양식의 엑셀 파일 취합 (결과 저장 폴더, 작업 대상 폴더) 
def C_xl(TPath, Result):

    # 작업할 파일이 모여있는 대상폴더를 선택하고 해당 폴더안의 엑셀파일 리스트 읽어오기
    xlfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]xls', f)]

    # Data를 취합할 DataFrame 생성해두기 (데이터 없는 Blank 구조)
    df = pd.DataFrame()

    # 파일 이름 읽어오면서 값 적재하는 로직 반복하기
    for xlfile in xlfiles:
        dfxl = pd.read_excel(os.path.join(TPath,xlfile))
        df = pd.concat([df,dfxl])

    # 결과파일 저장하기
    df.to_excel(os.path.join(Result, 'Result.xlsx'), index=False)

# 유사한 양식의 엑셀 파일 취합 (결과 저장 폴더, 작업 대상 폴더, 양식파일) 
def M_xl(TPath, RPath, Format):
        
    # Master 정보 읽어오기 (dfM : Mapping 정보, dfH : Head 정보)
    dfM = pd.read_excel(Format, sheet_name = 'Mapping')
    dfH = pd.read_excel(Format, sheet_name = 'Head')

    # Target Data가 모여 있는 폴더의 파일리스트 가져오기
    xlfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]xls', f)]

    # 결합할 DataFrame 초기값 설정하기 
    df = pd.DataFrame()

    # 엑셀 파일 리스트에 따라 모든 엑셀 파일을 읽어올때 까지 반복문
    for xlfile in xlfiles:
        # Target Data 엑셀 파일 정보 불러와 dfW라는 이름의 DataFrame에 저장하기
        dfW = pd.read_excel(os.path.join(TPath,xlfile))
        dfW['소속'] = xlfile[:xlfile.find('.')]

        # Mapping Table 정보를 모두 읽을때 까지 반복하기
        for index, row in dfM.iterrows():
            # dfW 칼럼 이름 변경하기
            dfW.rename(columns = {row[0]:row[1]}, inplace = True)
        # 엑셀 파일 결합하기
        df = pd.concat([df,dfW])

    # 결과 정렬하여 보기 (dfA : 소속 정보를 추가하기 위한 임시 DataFrame, 
    #                   dfR : 최종 결과 DataFrame)
    dfA = pd.DataFrame(['소속'], columns = ['Header'])
    dfH = pd.concat([dfH, dfA], ignore_index=True)
    dfR = pd.DataFrame(df, columns = dfH['Header'])

    # 결과 저장하기
    dfR.to_excel(os.path.join(RPath,'Result.xlsx'), index=False)

# 파워포인트 폰트 한번에 통일 (결과 저장폴더, 작업 대상 폴더, 영문폰트, 한글폰트) 
def PptFont (TPath, RPath, EFont, HFont):

    # 파워포인트 어플리케이션 개체 만들기
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")

    # 읽어올 폴더에서 파일 이름 리스트 가져오기
    pptfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]ppt', f)]

    # 파일에 대해 반복하기
    for pptfile in pptfiles:

        # 파워포인트에서 파일 읽어오기
        ppt = powerpoint.Presentations.Open(os.path.join(TPath,pptfile))

        # 모든 slide에 대해 반복
        for slide in ppt.Slides:
            # 모든 shape에 대해 반복
            for shape in slide.shapes:
                # TextFrame이 있는 경우
                if shape.HasTextFrame == -1:
                    shape.TextFrame.TextRange.Font.NameFarEast = HFont
                    shape.TextFrame.TextRange.Font.Name = EFont
                # 표가 있는 경우
                elif shape.HasTable == -1:
                    # 모든 행에 대하여 반복
                    for row in shape.Table.Rows:
                        # 행을 구성하는 모든 셀에 대하여 반복
                        for cell in row.cells:
                            cell.Shape.TextFrame.TextRange.Font.NameFarEast = HFont
                            cell.Shape.TextFrame.TextRange.Font.Name = EFont
                
                # 그 외의 경우
                else:
                    # try 아래 구문은 시도하다가 에러가 발생하거나 실패할 경우 except로 이동
                    try:
                        # Group을 구성하는 모든 구성 Shape에 대해 반복. 
                        # 개체가 Group이 아닐때 에러
                        for GI in shape.GroupItems:
                            if GI.HasTextFrame == -1:
                                GI.TextFrame.TextRange.Font.NameFarEast = HFont
                                GI.TextFrame.TextRange.Font.Name = EFont
                    # 에러 발생시 다음 단계로 이동
                    except:
                        pass


        # 파워포인트에서 파일 저장하고 파워포인트 닫기
        ppt.SaveAs (RPath + '\\R_' + pptfile)
        ppt.Close ()

# PPT 단어 일괄 수정 (결과 저장 폴더, 작업 대상 폴더, 양식파일) 
def PPT_Dic(TPath, RPath, DPath):

    # 읽어올 폴더에서 파일 이름 리스트 가져오기
    pptfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]ppt', f)]

    # 사전 정보 읽어오기
    df = pd.read_excel(DPath)

    # 파일에 대해 반복하기
    for pptfile in pptfiles:
        
        # 파워포인트에서 파일 읽어오기
        prs = Presentation(os.path.join(TPath,pptfile))
        for slide in prs.slides:
                for shape in slide.shapes:
                    # shape 단에서 다룰 내용이 있으면, shape 관련 내용 수정

                    # shape에 Text Frame이 포함되어 있는 경우
                    if shape.has_text_frame:
                            # text frame이 포함된 shape에서만 다룰 내용이 있으면 관련 내용 수정
                            para_function (shape,df)

                    # shape가 표인 경우
                    if shape.has_table:
                            # table이 포함된 shape에서만 다룰 내용이 있으면 관련 내용 수정

                            for row in shape.table.rows:
                                # 표의 행 단위로 다룰 내용이 있으면 관련 내용 수정

                                for cell in row.cells:
                                        # 표 안의 각 칸 단위로 다룰 내용이 있으면 관련 내용 수정
                                        para_function (cell,df)

                    # shape가 Group으로 묶인 경우
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:

                            for shp in shape.shapes:
                                if shp.has_text_frame:
                                        # text frame이 포함된 shape에서만 다룰 내용이 있으면 관련 내용 수정
                                        para_function (shp,df)
        # 결과 저장                                    
        prs.save(os.path.join(RPath, 'R_' + pptfile))

# 문단 단위 글자 변경 sub 함수 
def para_function (shp,df):
    for paragraph in shp.text_frame.paragraphs:
        for run in paragraph.runs:
            # 단어장 정보 찾아서 반복
            for index, row in df.iterrows():
                run.text = run.text.replace(row['Old'],row['New'])

# 워드 단어 일괄 수정 (결과 저장 폴더, 작업 대상 폴더, 양식파일) 
def Doc_Dic(TPath, RPath, DPath):
    # 읽어올 폴더에서 파일 이름 리스트 가져오기
    docfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]doc', f)]

    # 사전 정보 읽어오기
    df = pd.read_excel(DPath)

    for docfile in docfiles:
        # 워드 파일 읽어오기
        doc = Document(os.path.join(TPath,docfile))

        # 찾을 단어에 대해서 반복
        for index, dfrow in df.iterrows():
            # 각 문단에서 반복문 실행
            for paragraph in doc.paragraphs:
                # 문단의 각 run(띄어쓰기) 단위로 반복문 실행
                for run in paragraph.runs:
                    # 찾을 단어가 있을 때만 실행    
                    if run.text.find(dfrow['Old']) != -1 :
                        # 원하는 단어 찾아 바꾸기
                        run.text = run.text.replace(dfrow['Old'],dfrow['New'])

            # 각 표에서 반복문 실행
            for table in doc.tables:
                # 표의 각 행단위로 반복문 실행
                for row in table.rows:
                    # 행 내부의 Cell 단위로 반복문 실행
                    for cell in row.cells:
                        # 원하는 단어 찾아 바꾸기
                        cell.text = cell.text.replace(dfrow['Old'],dfrow['New'])

        doc.save(os.path.join(RPath,docfile))

# 폴더내 문서 PDF 파일 저장 (결과 저장 폴더, 작업 대상 폴더) 
def D_PDF(TPath,RPath):
        
    # 읽어올 폴더에서 파일 이름 리스트 가져오기
    pptfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]ppt', f)]
    xlfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]xls', f)]
    wordfiles = [f for f in os.listdir(TPath) if re.match(r'.*[.]doc', f)]

    # PPT 파일 PDF 전환하기
    # 파워포인트 어플리케이션 개체 만들기
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")

    # 파일 이름별 반복하기
    for pptfile in pptfiles:
        # 파워포인트에서 파일 읽어오기
        prs = powerpoint.Presentations.Open(os.path.join(TPath, pptfile))

        # PDF 파일로 저장하기 ※ 파일형태 PDF를 의미하는 번호 : 32
        savefile = RPath + '/ppt_' + pptfile[:pptfile.find('.')] + '.pdf'
        savefile = savefile.replace('/','\\')
        prs.SaveAs(savefile, 32)

        # 읽었던 파일 닫기
        prs.Close()

    # 파워포인트 어플리케이션 종료
    powerpoint.Quit()


    # 워드 파일 PDF 전환하기
    # 워드 어플리케이션 개체 만들기
    word = win32com.client.Dispatch("Word.Application")

    # 워드 파일이름별 반복하기
    for wordfile in wordfiles:
        # 워드에서 파일 읽어오기
        doc = word.Documents.Open(os.path.join(TPath,wordfile))

        # PDF 파일로 저장하기 ※ Word 파일형태 PDF를 의미하는 번호 : 17
        savefile = 'doc_' + wordfile[:wordfile.find('.')] + '.pdf'
        savefile = os.path.join(RPath, savefile)
        savefile = savefile.replace('/','\\')
        doc.SaveAs(savefile, 17)

        # 읽었던 파일 닫기
        doc.Close()

    # 워드 어플리케이션 종료
    word.Quit()

    # 엑셀 파일 PDF 전환하기
    # 엑셀 어플리케이션 개체 만들기
    excel = win32com.client.Dispatch("Excel.Application")

    # 엑셀 파일 이름별 반복하기
    for xlfile in xlfiles:
        # 엑셀에서 파일 읽어오기
        wb = excel.Workbooks.Open(os.path.join(TPath,xlfile))

        # PDF 파일로 저장하기 ※ 파일형태 PDF를 의미하는 번호 : 32
        savefile = 'xl_' + xlfile[:xlfile.find('.')] + '.pdf'
        wb.ExportAsFixedFormat(0, os.path.join(RPath,savefile))

        # 읽었던 파일 닫기
        wb.Close()

    # 엑셀 어플리케이션 종료
    excel.Quit()

# 폴더내 이미지 모아 PDF 파일 저장 (결과 저장 폴더, 작업 대상 폴더)
def Img_PDF(TPath,RPath):
    # 가로방향 PDF 파일 만들기 
    pdf = FPDF('L')

    # Image 파일 리스트 불러오기, 확장자명에 jpg, png, gif가 포함된 파일 전체
    files = [f for f in os.listdir(TPath) if re.match('.*([.]jpg|[.]png|[.]gif)', f)]

    # 파일 이름별로 반복
    for file in files:
        # img를 PDF로 바꿔서 페이지 추가
        pdf.add_page()
        # 이미지 불러와서 이미지 크기 확인하기
        img = Image.open(os.path.join(TPath,file))
        width, height = img.size
        # 이미지 크기를 mm 단위로 환산 (1px = 0.264583)
        width, height = float(width * 0.264583), float(height*0.264583)
        # 폭과 높이 A4 비율과 비교해서 크기 조정하기
        wid = 297 if width / height >= 297 / 210 else width * 210 / height 
        hei = height * 297 / width if width / height >= 297 / 210 else 210
        # wid, hei라는 이름으로 구한 치수 기준으로 이미지 삽입하기
        pdf.image(os.path.join(TPath,file), x= (297 - wid) / 2, y= (210 - hei) / 2, w = wid)

    # PDF 파일 저장하기
    pdf.output(os.path.join(RPath,"IMG2PDF.pdf"), "F")

# 폴더내 PDF 파일 암호 설정 (결과 저장 폴더, 작업 대상 폴더, 암호) 
def pw_PDF(TPath,RPath,pw):

    # Image 파일 리스트 불러오기, 확장자명에 jpg, png, gif가 포함된 파일 전체
    pdffiles = [f for f in os.listdir(TPath) if re.match('.*([.]pdf)', f)]

    # pdf 파일 리스트에 대해서 반복문
    for pdffile in pdffiles:
        # 파일을 열고 Password 설정 ※ rb 옵션은 파일을 읽기 위한 binary 포멧이라는 의미
        input_pdf = PdfReader(open(os.path.join(TPath,pdffile), "rb"))

        # PDF 파일 저장을 위한 Writer 생성하고 패스워드 설정하여 내용 입력
        output_pdf = PdfWriter()
        output_pdf.append_pages_from_reader(input_pdf)
        output_pdf.encrypt(pw)

        # 만들어진 결과 파일을 저장 ※ wb 옵션은 파일을 저장하기 위한 binary 포멧이라는 의미
        output_pdf.write(open(os.path.join(RPath, 'pw_' + pdffile), "wb"))        
